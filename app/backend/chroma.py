import chromadb
from chromadb.config import Settings
from langchain_chroma import Chroma
from langchain_google_genai import GoogleGenerativeAIEmbeddings
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.docstore.document import Document
from PyPDF2 import PdfReader
from PIL import Image
from docx import Document as DocxDocument
from pptx import Presentation
import os
import logging
import pytesseract

from dotenv import load_dotenv

load_dotenv()

CHROMA_DIR = "chroma_db"

# Inisialisasi ChromaDB
chroma_client = chromadb.Client(Settings(persist_directory=CHROMA_DIR))

# Embedding Google GenAI
embedding = GoogleGenerativeAIEmbeddings(
    model="models/text-embedding-004",
    google_api_key=os.getenv("GOOGLE_API_KEY", "")
    )

# LangChain VectorStore
vectorstore = Chroma(
    client=chroma_client,
    collection_name="rag_docs",
    embedding_function=embedding,
)

logging.basicConfig(level=logging.INFO)

def extract_text_from_pdf(file_path):
    reader = PdfReader(file_path)
    text = ""
    for page in reader.pages:
        text += page.extract_text() or ""
    return text

def extract_text_from_image(file_path):
    image = Image.open(file_path)
    text = pytesseract.image_to_string(image, lang='eng')
    return text

def extract_text_from_docx(file_path):
    doc = DocxDocument(file_path)
    text = "\n".join([para.text for para in doc.paragraphs])
    return text

def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def index_file(file_path, filetype, metadata=None):
    logging.info(f"[ChromaDB] Mulai indexing file: {file_path} ({filetype}) dengan metadata: {metadata}")
    if filetype == "application/pdf":
        text = extract_text_from_pdf(file_path)
    elif filetype.startswith("image/"):
        text = extract_text_from_image(file_path)
    elif filetype == "application/vnd.openxmlformats-officedocument.wordprocessingml.document" or file_path.endswith(".docx"):
        text = extract_text_from_docx(file_path)
    elif filetype == "application/vnd.openxmlformats-officedocument.presentationml.presentation" or file_path.endswith(".pptx"):
        text = extract_text_from_pptx(file_path)
    else:
        logging.warning(f"[ChromaDB] Filetype tidak didukung: {filetype}")
        return False
    splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=100)
    chunks = splitter.split_text(text)
    logging.info(f"[ChromaDB] Jumlah chunk yang akan di-embedding: {len(chunks)}")
    docs = [Document(page_content=chunk, metadata=metadata or {}) for chunk in chunks]
    # Logging sebelum embedding
    logging.info(f"[ChromaDB] Mulai proses embedding dengan Google GenAI untuk {len(docs)} dokumen...")
    vectorstore.add_documents(docs)
    logging.info(f"[ChromaDB] Selesai indexing dan embedding file: {file_path}")
    return True

def query_rag(query, top_k=3, file_path=None):
    if file_path:
        from langchain.text_splitter import RecursiveCharacterTextSplitter
        from langchain.docstore.document import Document
        if file_path.endswith('.pdf'):
            text = extract_text_from_pdf(file_path)
        elif file_path.endswith('.docx'):
            text = extract_text_from_docx(file_path)
        elif file_path.endswith('.pptx'):
            text = extract_text_from_pptx(file_path)
        elif file_path.lower().endswith(('.png', '.jpg', '.jpeg', '.bmp', '.gif')):
            text = extract_text_from_image(file_path)
        else:
            text = ""
        splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=100)
        docs = [Document(page_content=chunk) for chunk in splitter.split_text(text)]
        results = [(doc, 0) for doc in docs if query.lower() in doc.page_content.lower()][:top_k]
        return [doc for doc, _ in results]
    else:
        docs = vectorstore.similarity_search(query, k=top_k)
        return docs
