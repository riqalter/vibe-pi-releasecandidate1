import chromadb
from chromadb.config import Settings
from langchain_chroma import Chroma
from langchain_google_genai import GoogleGenerativeAIEmbeddings
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.docstore.document import Document
from PyPDF2 import PdfReader
from PIL import Image
from docx import Document as DocxDocument
from pptx import Presentation
import os
import logging
import pytesseract
import uuid

from dotenv import load_dotenv

load_dotenv()

CHROMA_DIR = "chroma_db"

chroma_client = chromadb.Client(Settings(persist_directory=CHROMA_DIR))

embedding = GoogleGenerativeAIEmbeddings(
    model="models/text-embedding-004",
    GOOGLE_API_KEY=os.getenv("GOOGLE_API_KEY", "")
    )

vectorstore = Chroma(
    client=chroma_client,
    collection_name="rag_docs",
    embedding_function=embedding,
)

logging.basicConfig(level=logging.INFO)

def extract_text_from_pdf(file_path):
    reader = PdfReader(file_path)
    docs = []
    for i, page in enumerate(reader.pages):
        content = page.extract_text() or ""
        if content:
            docs.append(Document(page_content=content, metadata={"page_number": i + 1}))
    return docs

def extract_text_from_image(file_path):
    image = Image.open(file_path)
    text = pytesseract.image_to_string(image, lang='eng')
    if text:
        return [Document(page_content=text, metadata={"page_number": 1})]
    return []

def extract_text_from_docx(file_path):
    doc = DocxDocument(file_path)
    docs = []
    for i, para in enumerate(doc.paragraphs):
        if para.text:
            docs.append(Document(page_content=para.text, metadata={"paragraph_number": i + 1}))
    return docs

def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    docs = []
    for i, slide in enumerate(prs.slides):
        slide_text = ""
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text += shape.text + "\n"
        if slide_text:
            docs.append(Document(page_content=slide_text, metadata={"slide_number": i + 1}))
    return docs

def index_file(file_path, filetype, metadata=None):
    logging.info(f"[ChromaDB] Mulai indexing file: {file_path} ({filetype}) dengan metadata: {metadata}")
    docs = []
    if filetype == "application/pdf":
        docs = extract_text_from_pdf(file_path)
    elif filetype.startswith("image/"):
        docs = extract_text_from_image(file_path)
    elif filetype == "application/vnd.openxmlformats-officedocument.wordprocessingprocessingml.document" or file_path.endswith(".docx"):
        docs = extract_text_from_docx(file_path)
    elif filetype == "application/vnd.openxmlformats-officedocument.presentationml.presentation" or file_path.endswith(".pptx"):
        docs = extract_text_from_pptx(file_path)
    else:
        logging.warning(f"[ChromaDB] Filetype tidak didukung: {filetype}")
        return False

    if not docs:
        logging.warning(f"[ChromaDB] Tidak ada teks yang diekstrak dari file: {file_path}")
        return False

    for doc in docs:
        doc.metadata.update(metadata or {})

    splitter = RecursiveCharacterTextSplitter(chunk_size=2000, chunk_overlap=200)
    chunks = splitter.split_documents(docs)
    
    logging.info(f"[ChromaDB] Jumlah chunk yang akan di-embedding: {len(chunks)}")
    if chunks:
        vectorstore.add_documents(chunks)
        logging.info(f"[ChromaDB] Selesai indexing dan embedding file: {file_path}")
        return True
    else:
        logging.warning(f"[ChromaDB] Tidak ada chunk yang dihasilkan setelah splitting: {file_path}")
        return False

def query_rag(query, top_k=10, file_path=None):
    if file_path:
        from langchain.text_splitter import RecursiveCharacterTextSplitter
        from langchain.docstore.document import Document
        if file_path.endswith('.pdf'):
            loaded_docs = extract_text_from_pdf(file_path)
        elif file_path.endswith('.docx'):
            loaded_docs = extract_text_from_docx(file_path)
        elif file_path.endswith('.pptx'):
            loaded_docs = extract_text_from_pptx(file_path)
        elif file_path.lower().endswith(('.png', '.jpg', '.jpeg', '.bmp', '.gif')):
            loaded_docs = extract_text_from_image(file_path)
        else:
            loaded_docs = []

        filename = os.path.basename(file_path)
        for doc in loaded_docs:
            doc.metadata['filename'] = filename

        splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=100)
        chunks = splitter.split_documents(loaded_docs)
        logging.info(f"[ChromaDB] Chunks metadata before adding to temp_collection: {[c.metadata for c in chunks]}")
        
        temp_collection_name = f"temp_rag_docs_{uuid.uuid4().hex}"
        temp_collection = chroma_client.create_collection(name=temp_collection_name)
        
        temp_collection.add(
            documents=[chunk.page_content for chunk in chunks],
            metadatas=[chunk.metadata for chunk in chunks],
            ids=[f"doc_{i}" for i in range(len(chunks))]
        )
        
        results = temp_collection.query(
            query_texts=[query],
            n_results=top_k,
            include=['documents', 'metadatas']
        )
        
        docs = []
        if results and results['documents']:
            for i, doc_content in enumerate(results['documents'][0]):
                doc_metadata = results['metadatas'][0][i]
                docs.append(Document(page_content=doc_content, metadata=doc_metadata))
        logging.info(f"[ChromaDB] Docs metadata after retrieval from temp_collection: {[d.metadata for d in docs]}")
        
        return docs
    else:
        docs = vectorstore.similarity_search(query, k=top_k)
        return docs
